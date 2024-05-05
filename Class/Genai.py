import os
import hashlib
from pathlib import Path
import google.generativeai as genai
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


class GenerativeAI:
    def __init__(self, api_key):
        self.api_key = api_key
        genai.configure(api_key=self.api_key)

    def questionFileCode(self, text, file):
        generation_config = {
            "temperature": 1,
            "top_p": 0.95,
            "top_k": 0,
            "max_output_tokens": 8192,
        }

        safety_settings = [
            {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        ]

        model = genai.GenerativeModel(
            model_name="gemini-1.5-pro-latest",
            generation_config=generation_config,
            safety_settings=safety_settings,
        )
        uploaded_files = []

        def upload_if_needed(pathname: str) -> list[str]:
            path = Path(pathname)
            hash_id = hashlib.sha256(path.read_bytes()).hexdigest()
            try:
                existing_file = genai.get_file(name=hash_id)
                return [existing_file.uri]
            except:
                pass
            uploaded_files.append(genai.upload_file(path=path, display_name=hash_id))
            return [uploaded_files[-1].uri]

        def read_text_file(pathname: str) -> str:
            with open(pathname, "r", encoding="utf-8") as file:
                content = file.read()
            return content

        convo = model.start_chat(
            history=[
                {
                    "role": "user",
                    "parts": [read_text_file(file)]
                },
                {
                    "role": "model",
                    "parts": ["Ok, recebi o arquivo."]
                },
            ]
        )

        convo.send_message(text)
        response = convo.last.text

        for uploaded_file in uploaded_files:
            genai.delete_file(name=uploaded_file.name)

        return response
    
    def question(self, text):  # Added self as the first parameter
        generation_config = {
        "temperature": 1,
        "top_p": 0.95,
        "top_k": 0,
        "max_output_tokens": 8192,
        }

        safety_settings = [
        {
            "category": "HARM_CATEGORY_HARASSMENT",
            "threshold": "BLOCK_MEDIUM_AND_ABOVE"
        },
        {
            "category": "HARM_CATEGORY_HATE_SPEECH",
            "threshold": "BLOCK_MEDIUM_AND_ABOVE"
        },
        {
            "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
            "threshold": "BLOCK_MEDIUM_AND_ABOVE"
        },
        {
            "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
            "threshold": "BLOCK_MEDIUM_AND_ABOVE"
        },
        ]

        model = genai.GenerativeModel(model_name="gemini-1.5-pro-latest",
                                    generation_config=generation_config,
                                    safety_settings=safety_settings)

        convo = model.start_chat(history=[
        {
            "role": "user",
            "parts": ["aparitir de agora o seu nome sera helena, um assiste virtual de estudo para programadores e entusiasta da tecnologia cirado por walter alexandre santana(programador full static web, javascript, python e c#, residente em angola, africa estudante da escola 42 luanda) em 23 de marco de 2024"]
        },
        {
            "role": "model",
            "parts": ["Olá! A partir de agora, pode me chamar de Helena. Fui criada por Walter Alexandre Santana para ser sua assistente virtual. Em que posso ajudar hoje?"]
        },
        {
            "role": "user",
            "parts": ["Pare de usar emojis nas suas respostas"]  # Corrected spelling of "emojis"
        },
        {
            "role": "model",
            "parts": ["Ok!"]
        },
        ])
        
        convo.send_message(text)
        return convo.last.text
    
    def questionFile(self, text, file_type, file_path):  # Added self as the first parameter
        generation_config = {
            "temperature": 1,
            "top_p": 0.95,
            "top_k": 0,
            "max_output_tokens": 8192,
        }

        safety_settings = [
            {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        ]

        model = genai.GenerativeModel(
            model_name="gemini-1.5-pro-latest",
            generation_config=generation_config,
            safety_settings=safety_settings,
        )

        uploaded_files = []

        def upload_if_needed(pathname: str) -> list[str]:
            path = Path(pathname)
            hash_id = hashlib.sha256(path.read_bytes()).hexdigest()
            try:
                existing_file = genai.get_file(name=hash_id)
                return [existing_file.uri]
            except Exception as e:
                print(f"Error uploading file: {e}")
                pass
            uploaded_files.append(genai.upload_file(path=path, display_name=hash_id))
            return [uploaded_files[-1].uri]

        def read_docx_file(pathname: str) -> str:
            try:
                doc = Document(pathname)
                content = ""
                for paragraph in doc.paragraphs:
                    content += paragraph.text + "\n"
                return content
            except Exception as e:
                print(f"Error reading Word document: {e}")
                return ""

        def read_xlsx_file(pathname: str) -> str:
            try:
                wb = load_workbook(pathname)
                content = ""
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        content += "\t".join(str(cell) for cell in row) + "\n"
                return content
            except Exception as e:
                print(f"Error reading Excel file: {e}")
                return ""

        def read_pptx_file(pathname: str) -> str:
            try:
                prs = Presentation(pathname)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content
            except Exception as e:
                print(f"Error reading PowerPoint file: {e}")
                return ""

        if file_type == 'excel':
            content_reader = read_xlsx_file
        elif file_type == 'word':
            content_reader = read_docx_file
        elif file_type == 'powerpoint':
            content_reader = read_pptx_file
        else:
            print("Unsupported file type.")
            return

        try:
            convo = model.start_chat(
                history=[
                    {
                        "role": "user",
                        "parts": ["aparitir de agora o seu nome sera helena, um assiste virtual de estudo para programadores e entusiasta da tecnologia cirado por walter alexandre santana(programador full static web, javascript, python e c#, residente em angola, africa estudante da escola 42 luanda) em 23 de marco de 2024"]
                    },
                    {
                        "role": "model",
                        "parts": ["Olá! A partir de agora, pode me chamar de Helena. Fui criada por Walter Alexandre Santana para ser sua assistente virtual. Em que posso ajudar hoje?"]
                    },
                    {
                        "role": "user",
                        "parts": ["Pare de usar emojis nas suas respostas"]
                    },
                    {
                        "role": "model",
                        "parts": ["Ok!"]
                    },
                    {
                        "role": "user",
                        "parts": [content_reader(file_path)]
                    },
                    {
                        "role": "model",
                        "parts": ["Ok, received the file."]
                    },
                ]
            )

            convo.send_message(text)
            return convo.last.text
        except Exception as e:
            return f"Error communicating with Generative AI model: {e}"

        for uploaded_file in uploaded_files:
            try:
                genai.delete_file(name=uploaded_file.name)
            except Exception as e:
                return f"Error deleting file: {e}"
            
    def questionFilePDF(self, text, file):  # Added self as the first parameter

        generation_config = {
            "temperature": 1,
            "top_p": 0.95,
            "top_k": 0,
            "max_output_tokens": 8192,
        }

        safety_settings = [
            {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        ]

        model = genai.GenerativeModel(
            model_name="gemini-1.5-pro-latest",
            generation_config=generation_config,
            safety_settings=safety_settings,
        )

        uploaded_files = []

        def upload_if_needed(pathname: str) -> list[str]:
            path = Path(pathname)
            hash_id = hashlib.sha256(path.read_bytes()).hexdigest()
            try:
                existing_file = genai.get_file(name=hash_id)
                return [existing_file.uri]
            except:
                pass
            uploaded_files.append(genai.upload_file(path=path, display_name=hash_id))
            return [uploaded_files[-1].uri]

        def extract_pdf_pages(pathname: str) -> list[str]:
            parts = []
            with open(pathname, "rb") as file:
                pdf_reader = PyPDF2.PdfFileReader(file)
                for page_num in range(pdf_reader.numPages):
                    page = pdf_reader.getPage(page_num).extractText()
                    parts.append(page)
            return parts

        convo = model.start_chat(
            history=[
                {
                    "role": "user",
                    "parts": extract_pdf_pages(file)
                },
                {
                    "role": "model",
                    "parts": ["Ok, recebi o arquivo."]
                },
            ]
        )

        convo.send_message(text)
        response = convo.last.text

        for uploaded_file in uploaded_files:
            genai.delete_file(name=uploaded_file.name)

        return response
